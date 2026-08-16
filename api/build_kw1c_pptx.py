#!/usr/bin/env python3
"""
build_kw1c_pptx.py - Genereer een KW1C-huisstijl PowerPoint op basis van content.json.

Gebruik:
    python build_kw1c_pptx.py content.json output.pptx
"""

import argparse
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("python-pptx niet gevonden. Installeer met: pip install python-pptx --break-system-packages")
    sys.exit(1)

# ── KW1C huisstijl kleuren ────────────────────────────────────────────────────
KW1C_BLAUW  = RGBColor(0x00, 0x9B, 0xD2)
KW1C_DONKER = RGBColor(0x00, 0x85, 0xCA)
KW1C_ROOD   = RGBColor(0xFF, 0x00, 0x00)
KW1C_WIT    = RGBColor(0xFF, 0xFF, 0xFF)
KW1C_LICHT  = RGBColor(0xFA, 0xF8, 0xF6)

# ── Layout-nummers ────────────────────────────────────────────────────────────
LAYOUT = {
    "titelslide":      0,
    "inhoudsopgave":   1,
    "blauw":           2,
    "rood":            3,
    "quote":           4,
    "sectie":          5,
    "tweekolommen":    6,
    "afbeelding":      7,
    "afbeelding_only": 8,
    "afsluiting":      9,
}

# Contactgegevens die op de titelslide komen (bewerkbaar op de slide zelf)
CONTACT_TEKST = "Marja Müskens  ★  m.muskens@kw1c.nl  ★  06 48 267 988"

# Rode balk op titelslide: y=17.71cm, h=1.34cm, volledige breedte 33.87cm
RODE_BALK_Y   = 17.71
RODE_BALK_H   = 1.34
SLIDE_BREEDTE = 33.87


def get_layout(prs, name):
    idx = LAYOUT.get(name, 2)
    return prs.slide_layouts[idx]


def add_slide(prs, layout_name):
    return prs.slides.add_slide(get_layout(prs, layout_name))


def set_ph(slide, ph_idx, text, bold=None, font_size=None, color=None, no_bullet=False):
    """Vul een placeholder. Optioneel bullet uitzetten."""
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if no_bullet:
        pPr = p._p.get_or_add_pPr()
        # Verwijder eventuele bestaande buNone
        for old in pPr.findall(qn('a:buNone')):
            pPr.remove(old)
        etree.SubElement(pPr, qn('a:buNone'))
    run = p.add_run()
    run.text = text
    if bold is not None:
        run.font.bold = bold
    if font_size is not None:
        run.font.size = Pt(font_size)
    if color is not None:
        run.font.color.rgb = color


def set_bullets(slide, ph_idx, bullets, font_size=20):
    """Vul een body-placeholder met bulletlijst."""
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return
    tf = ph.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)


def add_textbox(slide, text, x_cm, y_cm, w_cm, h_cm,
                font_size=20, color=None, bold=False,
                align=PP_ALIGN.CENTER):
    """Vrije tekstbox op de slide (niet in diamodel - direct bewerkbaar in PowerPoint)."""
    txBox = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = color
    return txBox


# ── Slide builders ────────────────────────────────────────────────────────────

def build_titelslide(prs, data):
    slide = add_slide(prs, "titelslide")
    title = data.get("titel", "TITEL")
    subtitle = data.get("ondertitel", "")
    set_ph(slide, 0, title.upper(), bold=False, font_size=48, color=KW1C_WIT)
    if subtitle:
        set_ph(slide, 1, subtitle, font_size=24, color=KW1C_WIT)

    # Contactgegevens als vrije tekstbox bovenop de rode balk
    # Gecentreerd, wit, 20pt - direct klikbaar en aanpasbaar in PowerPoint
    contact = data.get("contact", CONTACT_TEKST)
    if contact:
        # Kleine marge aan de zijkanten (1cm), verticaal gecentreerd in de balk
        add_textbox(
            slide, contact,
            x_cm=1.0, y_cm=RODE_BALK_Y,
            w_cm=SLIDE_BREEDTE - 2.0, h_cm=RODE_BALK_H,
            font_size=20, color=KW1C_WIT, bold=False,
            align=PP_ALIGN.CENTER
        )
    return slide


def build_inhoudsopgave(prs, data):
    slide = add_slide(prs, "inhoudsopgave")
    set_ph(slide, 0, data.get("titel", "INHOUDSOPGAVE").upper(), bold=False)
    items = data.get("items", [])
    if items:
        set_bullets(slide, 1, items, font_size=20)
    return slide


def build_contentslide(prs, data):
    variant = data.get("variant", "blauw")
    layout_name = "rood" if variant == "rood" else "blauw"
    slide = add_slide(prs, layout_name)
    set_ph(slide, 0, data.get("titel", "").upper(), bold=False)
    bullets = data.get("bullets", [])
    if bullets:
        set_bullets(slide, 1, bullets, font_size=20)
    elif data.get("tekst"):
        set_bullets(slide, 1, [data["tekst"]], font_size=20)
    return slide


def build_quote(prs, data):
    slide = add_slide(prs, "quote")
    set_ph(slide, 0, data.get("titel", "").upper(), bold=False)
    quote_text = data.get("quote", data.get("tekst", ""))
    # no_bullet=True: geen bulletpunt voor de quote
    set_ph(slide, 1, quote_text, font_size=28, color=KW1C_WIT, bold=False, no_bullet=True)
    bronvermelding = data.get("bron", "")
    if bronvermelding:
        set_ph(slide, 10, bronvermelding, font_size=14, color=KW1C_WIT, no_bullet=True)
    return slide


def build_sectie_intro(prs, data):
    slide = add_slide(prs, "sectie")
    set_ph(slide, 0, data.get("titel", "").upper(), bold=False)
    ondertitel = data.get("ondertitel", data.get("tekst", ""))
    if ondertitel:
        set_ph(slide, 1, ondertitel, font_size=22)
    return slide


def build_tweekolommen(prs, data):
    slide = add_slide(prs, "tweekolommen")
    set_ph(slide, 0, data.get("titel", "").upper(), bold=False)
    links = data.get("links", [])
    rechts = data.get("rechts", [])
    if links:
        set_bullets(slide, 1, links, font_size=20)
    if rechts:
        set_bullets(slide, 2, rechts, font_size=20)
    return slide


def build_afsluiting(prs, data):
    slide = add_slide(prs, "afsluiting")
    set_ph(slide, 0, data.get("titel", "VRAGEN?").upper(), bold=False)
    tekst = data.get("tekst", "")
    if tekst:
        items = [tekst] if isinstance(tekst, str) else tekst
        set_bullets(slide, 1, items, font_size=20)
    return slide


# ── Dispatcher ────────────────────────────────────────────────────────────────
BUILDERS = {
    "titelslide":    build_titelslide,
    "inhoudsopgave": build_inhoudsopgave,
    "content":       build_contentslide,
    "blauw":         build_contentslide,
    "rood":          build_contentslide,
    "quote":         build_quote,
    "sectie":        build_sectie_intro,
    "tweekolommen":  build_tweekolommen,
    "afsluiting":    build_afsluiting,
}


def build(content_path, output_path, template_path):
    with open(content_path, "r", encoding="utf-8") as f:
        content = json.load(f)

    prs = Presentation(template_path)

    # Verwijder voorbeeldslides uit de template
    slide_id_list = prs.slides._sldIdLst
    while len(slide_id_list) > 0:
        sld_id_elem = slide_id_list[0]
        r_id = sld_id_elem.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass
        slide_id_list.remove(sld_id_elem)

    slides_data = content.get("slides", [])
    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        builder = BUILDERS.get(slide_type, build_contentslide)
        builder(prs, slide_data)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(f"✓ Opgeslagen: {output_path} ({len(slides_data)} slides)")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("content")
    parser.add_argument("output")
    parser.add_argument("--template", default=None)
    args = parser.parse_args()

    if args.template:
        template_path = args.template
    else:
        script_dir = Path(__file__).parent
        template_path = script_dir.parent / "assets" / "kw1c_template.pptx"
        if not template_path.exists():
            print(f"Template niet gevonden op {template_path}. Geef --template mee.")
            sys.exit(1)

    build(args.content, args.output, str(template_path))


if __name__ == "__main__":
    main()
